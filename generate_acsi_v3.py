from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_slide_style(slide, title_text):
    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # Title Style
    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
                run.font.bold = True
                run.font.size = Pt(24)

def add_bullet_points(slide, points, font_size=Pt(16)):
    if not slide.placeholders[1]:
        return
    tf = slide.shapes.placeholders[1].text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = pt
        else:
            p = tf.add_paragraph()
            p.text = pt
        p.level = 0
        p.font.size = font_size
        p.font.color.rgb = RGBColor(0, 0, 0)

def create_presentation():
    prs = Presentation()

    # 1. Titre & Groupe
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_style(slide, "Projet SGRDMS")
    subtitle = slide.placeholders[1]
    subtitle.text = ("Système de Gestion des Rendez-vous et Dossiers Médicaux\n"
                     "Analyse et Conception de Système d'Information\n\n"
                     "Présenté par : [Membres du Groupe]\n"
                     "Licence 2 Génie Logiciel - Université Iba Der Thiam")
    for p in subtitle.text_frame.paragraphs:
        p.font.size = Pt(18)

    # 2. Contexte & Problématique
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "1. Contexte & Problématique")
    add_bullet_points(slide, [
        "Informatisation du centre de santé 'LE TROPICAL'.",
        "Transition du support papier vers un système numérique.",
        "Risques identifiés : Pertes de dossiers, erreurs de RDV.",
        "Problématiques : Lenteur de facturation, manque de traçabilité médicale.",
        "Besoin : Un système centralisé, sécurisé et multi-centres."
    ])

    # 3. Objectifs du Système
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "2. Objectifs du Système")
    add_bullet_points(slide, [
        "Centralisation des dossiers médicaux (Historique complet).",
        "Optimisation de la prise de RDV et gestion des files d'attente.",
        "Traçabilité des consultations, prescriptions et factures.",
        "Modules avancés : Téléconsultation, Gestion des Assurances.",
        "Sécurisation des accès (Rôles : Admin, Médecin, Patient, etc.)."
    ])

    # 4. Règles de Gestion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "3. Règles de Gestion (Extraits)")
    add_bullet_points(slide, [
        "RG1: Un patient possède un dossier médical unique.",
        "RG2: Un médecin appartient à un seul service spécialisé.",
        "RG5: Un médecin ne peut avoir deux RDV au même créneau horaire.",
        "RG13: Blocage automatique si interaction médicamenteuse à risque élevé.",
        "RG14: Gestion du tiers-payant (Part Patient vs Part Assurance)."
    ], font_size=Pt(14))

    # 5. Dictionnaire de Données
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "4. Dictionnaire de Données (Rigoureux)")
    
    rows, cols = 9, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.2), Inches(9.6), Inches(3.5)).table
    headers = ["Code", "Désignation", "Type", "Taille", "Nature"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        p = cell.text_frame.paragraphs[0]
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.bold = True
        p.font.size = Pt(11)

    dd_data = [
        ("ID_PAT", "Identifiant Patient", "AN", "10", "Élémentaire (PK)"),
        ("NOM_PAT", "Nom du Patient", "A", "50", "Élémentaire"),
        ("DAT_NAI", "Date de Naissance", "D", "8", "Élémentaire"),
        ("SEX_PAT", "Sexe (M/F)", "A", "1", "Domaine {M,F}"),
        ("ID_RDV", "Numéro de RDV", "AN", "10", "Élémentaire (PK)"),
        ("MT_TOTAL", "Montant Facture", "N", "12", "Calculé"),
        ("TX_ASSR", "Taux Prise en Charge", "N", "3", "Élémentaire"),
        ("STK_QTE", "Stock Disponible", "N", "6", "Élémentaire")
    ]
    for i, row in enumerate(dd_data, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    # 6. MCD - Vue d'ensemble
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_style(slide, "5. Modèle Conceptuel de Données (MCD)")
    img_path = "/home/mhd/sgrdms/mcd_image-1.png"
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.5), Inches(1.3), width=Inches(9))

    # 7. MCD - Détails des Entités
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "6. MCD : Focus Entités Majeures")
    
    rows, cols = 4, 2
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(3)).table
    table.cell(0,0).text = "Entité"
    table.cell(0,1).text = "Attributs Principaux"
    
    entities = [
        ("PATIENT", "id_patient, nom, prenom, date_naissance, sexe, groupe_sanguin, antecedents"),
        ("MEDECIN", "id_medecin, nom, specialite, telephone, teleconsultation_active"),
        ("CONSULTATION", "id_consultation, date, motif, diagnostic, tension, poids, taille")
    ]
    for i, (ent, attr) in enumerate(entities, start=1):
        table.cell(i, 0).text = ent
        table.cell(i, 1).text = attr
        table.cell(i, 1).text_frame.word_wrap = True
        for j in range(2):
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)

    # 8. MLD
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "7. Modèle Logique de Données (MLD)")
    mld_points = [
        "PATIENT (Id_Pat, Nom_Pat, Pre_Pat, Dat_Nai, Sex_Pat, Grp_San, ...)",
        "RENDEZ_VOUS (Id_Rdv, Dat_Rdv, Heu_Deb, Motif, #Id_Pat, #Id_Med)",
        "CONSULTATION (Id_Cons, Dat_Cons, Diagnostic, Poids, #Id_Rdv)",
        "SERVICE (Id_Ser, Nom_Ser, Description, #Id_Centre)",
        "LIGNE_ORD (Id_Lig, Dose, Frequence, #Id_Ord, #Id_Med)"
    ]
    add_bullet_points(slide, mld_points, font_size=Pt(14))

    # 9. Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_style(slide, "Conclusion")
    add_bullet_points(slide, [
        "Solution complète répondant aux besoins du centre 'LE TROPICAL'.",
        "Amélioration drastique de la qualité de prise en charge des patients.",
        "Sécurisation et pérennité des données médicales.",
        "Outil d'aide à la décision pour les médecins (Interactions).",
        "Perspective : Extension vers un réseau national de santé."
    ])

    final_path = "Presentation_SGRDMS_ACSI_V3.pptx"
    prs.save(final_path)
    print(f"Présentation générée : {final_path}")

if __name__ == "__main__":
    create_presentation()
