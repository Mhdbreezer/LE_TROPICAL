from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()

    # Colors
    MED_DARK_BLUE = RGBColor(44, 62, 80)   # #2c3e50
    MED_LIGHT_BLUE = RGBColor(52, 152, 219) # #3498db
    WHITE = RGBColor(255, 255, 255)

    def apply_header(slide, title_text):
        # Add a dark blue rectangle at the top
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = MED_DARK_BLUE
        shape.line.fill.background() # No border

        # Add title text manually to ensure styling
        title_shape = slide.shapes.title
        title_shape.text = title_text
        title_shape.left = Inches(0.5)
        title_shape.top = Inches(0.1)
        title_shape.width = prs.slide_width - Inches(1)
        title_shape.height = Inches(0.6)
        
        tf = title_shape.text_frame
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True

    def add_styled_slide(title_text, points):
        slide_layout = prs.slide_layouts[5] # Blank layout usually or Title Only
        slide = prs.slides.add_slide(slide_layout)
        apply_header(slide, title_text)

        # Add content area
        left = Inches(0.5)
        top = Inches(1.2)
        width = prs.slide_width - Inches(1)
        height = prs.slide_height - Inches(1.5)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        for point in points:
            p = tf.add_paragraph()
            p.text = "• " + point
            p.space_before = Pt(10)
            p.font.size = Pt(18)
            p.font.color.rgb = MED_DARK_BLUE

    # Slide 1: Title Slide (Special styling)
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
    # Background rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = MED_DARK_BLUE
    
    # Title
    title_box = slide.shapes.add_textbox(0, Inches(2), prs.slide_width, Inches(1.5))
    tf = title_box.text_frame
    tf.text = "SGRDMS"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.bold = True
    p.font.size = Pt(60)
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(0, Inches(3.5), prs.slide_width, Inches(2))
    tf = sub_box.text_frame
    tf.text = "Système de Gestion des Rendez-vous et\nDossiers Médicaux Spécialisés\n\nPrésenté par : Groupe SGRDMS\n2025-2026"
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(189, 195, 199) # Light gray

    # Slide 2: Introduction
    add_styled_slide("Introduction", [
        "Modernisation de la gestion hospitalière.",
        "Automatisation du suivi des patients et dossiers.",
        "Optimisation du flux de travail médical.",
        "Sécurisation des données sensibles.",
        "Accessibilité accrue via la téléconsultation."
    ])

    # Slide 3: Dictionnaire de Données
    add_styled_slide("Dictionnaire de Données", [
        "PATIENT : Informations civiles, contact, antécédents, assurance.",
        "MEDECIN : Spécialité, service, planning de disponibilité.",
        "CONSULTATION : Diagnostic, constantes vitales, observations.",
        "ORDONNANCE : Liste de médicaments, posologie, validité.",
        "UTILISATEUR : Identifiants sécurisés et rôles (RBAC)."
    ])

    # Slide 4: MCD (Conceptuel)
    add_styled_slide("Modèle Conceptuel (MCD)", [
        "Entités : Patient, Médecin, Service, Consultation, Pharmacie.",
        "Relations :",
        "  - Patient (1,N) -- Effectuer -- (1,1) RDV",
        "  - RDV (1,1) -- Aboutir -- (0,1) Consultation",
        "  - Consultation (1,1) -- Prescrire -- (0,N) Ordonnance",
        "  - Médicament (1,N) -- Composer -- (1,N) Ordonnance"
    ])

    # Slide 5: MLD (Logique)
    add_styled_slide("Modèle Logique (MLD)", [
        "Conversion en tables relationnelles :",
        "• Utilisateurs (id, login, password_hash, role)",
        "• Patients (id, nom, prenom, #assurance_id)",
        "• Medecins (id, nom, prenom, #service_id)",
        "• RendezVous (id, date, heure, #patient_id, #medecin_id)",
        "• LigneOrdonnance (id, #ord_id, #med_id, quantite)"
    ])

    # Slide 6: Architecture Technique
    add_styled_slide("Architecture & Technologies", [
        "Backend : Python 3.12 & Flask (Framework léger).",
        "Frontend : HTML5 / CSS3 / JavaScript (Bootstrap 5).",
        "Base de données : SQLAlchemy (ORM) & SQLite/PostgreSQL.",
        "Sécurité : Flask-Login & Bcrypt pour le hachage.",
        "Environnement : Isolation via Virtualenv."
    ])

    # Slide 7: Gestion des Rôles
    add_styled_slide("Gestion des Rôles (RBAC)", [
        "Administrateur : Contrôle total du système et des logs.",
        "Médecin : Accès aux dossiers médicaux et prescriptions.",
        "Pharmacien : Gestion du stock et validation des ordonnances.",
        "Réceptionniste : Gestion des flux physiques et RDV.",
        "Patient : Consultation de son propre historique."
    ])

    # Slide 8: Points Forts & Innovation
    add_styled_slide("Innovation & Atouts", [
        "Module de Téléconsultation (Visio & Chat).",
        "Gestion intelligente des stocks (Alertes de péremption).",
        "Vérification automatique des interactions médicamenteuses.",
        "File d'attente en temps réel pour la réception.",
        "Interface responsive adaptée aux tablettes et mobiles."
    ])

    # Slide 9: Démonstration Visuelle
    add_styled_slide("Interface Utilisateur", [
        "Tableaux de bord intuitifs par rôle.",
        "Formulaires de saisie rapide pour les médecins.",
        "Visualisation graphique des statistiques d'activité.",
        "Notifications automatiques pour les rappels de RDV.",
        "Recherche multicritères (Patient, Date, Médicament)."
    ])

    # Slide 10: Conclusion
    add_styled_slide("Conclusion", [
        "Solution intégrée, robuste et évolutive.",
        "Amélioration significative de la qualité des soins.",
        "Réduction des temps d'attente et erreurs de saisie.",
        "Prêt pour un déploiement en milieu clinique.",
        "Perspectives : IA pour l'aide au diagnostic."
    ])

    prs.save('Presentation_SGRDMS_Pro.pptx')
    print("Présentation 'Jolie' générée : Presentation_SGRDMS_Pro.pptx")

if __name__ == '__main__':
    create_presentation()


if __name__ == '__main__':
    create_presentation()
