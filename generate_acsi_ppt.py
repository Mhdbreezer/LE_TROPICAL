from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def set_slide_background_and_title_style(slide, title_text):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    if slide.shapes.title:
        title = slide.shapes.title
        title.text = title_text
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(0, 0, 0)
                run.font.bold = True
                run.font.size = Pt(28)

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_background_and_title_style(slide, "SGRDMS : Système de Gestion Médicale")
    subtitle = slide.placeholders[1]
    subtitle.text = "Analyse et Conception de Système d'Information (ACSI)\nDictionnaire de Données & Modélisation Merise\n\nPrésenté par : [Membres du Groupe]"

    # --- Slide 2: Contexte ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background_and_title_style(slide, "Contexte du Projet")
    tf = slide.shapes.placeholders[1].text_frame
    tf.text = "Le centre 'LE TROPICAL' souhaite automatiser :"
    for item in ["Gestion des dossiers patients", "Planification des rendez-vous", "Consultations et Ordonnances", "Pharmacie et Facturation"]:
        p = tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 1

    # --- Slide 3: Dictionnaire de Données (Rigoureux - Page 1) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background_and_title_style(slide, "Dictionnaire de Données (1/2)")
    
    rows, cols = 8, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.2), Inches(9.6), Inches(4)).table
    headers = ["Code", "Désignation", "Type", "Taille", "Contraintes"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        table.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True

    dd_data = [
        ("ID_PAT", "Identifiant Patient", "AN", "10", "Unique, Non Nul"),
        ("NOM_PAT", "Nom du Patient", "A", "50", "Majuscule"),
        ("PRE_PAT", "Prénom du Patient", "A", "50", "-"),
        ("DAT_NAI", "Date de Naissance", "D", "8", "JJ/MM/AAAA"),
        ("SEX_PAT", "Sexe (M/F)", "A", "1", "Domaine {M, F}"),
        ("TEL_PAT", "Téléphone", "N", "12", "Format International"),
        ("GRP_SAN", "Groupe Sanguin", "AN", "3", "{A+, A-, B+, ...}")
    ]
    for i, row in enumerate(dd_data, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)

    # --- Slide 4: Dictionnaire de Données (Rigoureux - Page 2) ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background_and_title_style(slide, "Dictionnaire de Données (2/2)")
    
    table = slide.shapes.add_table(rows, cols, Inches(0.2), Inches(1.2), Inches(9.6), Inches(4)).table
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        table.cell(0, i).text_frame.paragraphs[0].runs[0].font.bold = True

    dd_data_2 = [
        ("ID_RDV", "ID Rendez-vous", "AN", "10", "Unique"),
        ("DAT_RDV", "Date du RDV", "D", "8", ">= Date Jour"),
        ("MT_TOT", "Montant Total Facture", "N", "10", "> 0 (FCFA)"),
        ("TX_ASSR", "Taux Prise en Charge", "N", "3", "0 à 100%"),
        ("STK_MED", "Stock Médicament", "N", "5", ">= 0"),
        ("NIV_RISQ", "Niveau de Risque", "A", "10", "{Faible, Modéré, Elevé}"),
        ("ID_CONS", "ID Consultation", "AN", "10", "Unique")
    ]
    for i, row in enumerate(dd_data_2, start=1):
        for j, val in enumerate(row):
            table.cell(i, j).text = val
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(11)

    # --- Slide 5: Règles de Gestion ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background_and_title_style(slide, "Règles de Gestion Clés")
    tf = slide.shapes.placeholders[1].text_frame
    rgs = [
        "RG1: Un patient possède un dossier unique.",
        "RG3: Un RDV concerne 1 patient et 1 médecin (date/heure précise).",
        "RG5: Un médecin ne peut avoir deux RDV au même créneau.",
        "RG13: Blocage automatique si interaction médicamenteuse à risque élevé.",
        "RG14: Facturation avec prise en charge partielle par l'assurance."
    ]
    for rg in rgs:
        p = tf.add_paragraph()
        p.text = f"• {rg}"
        p.level = 0

    # --- Slide 6: MCD Image ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background_and_title_style(slide, "Modèle Conceptuel de Données (MCD)")
    mcd_img_path = "/home/mhd/sgrdms/mcd_image-1.png"
    if os.path.exists(mcd_img_path):
        slide.shapes.add_picture(mcd_img_path, Inches(0.2), Inches(1.2), width=Inches(9.6))

    # --- Slide 7: MLD ---
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    set_slide_background_and_title_style(slide, "Modèle Logique de Données (MLD)")
    tf = slide.shapes.placeholders[1].text_frame
    mld_text = (
        "PATIENT (Id_Pat, Nom_Pat, Pre_Pat, Dat_Nai, Sex_Pat, Tel_Pat, ...)\n"
        "RENDEZ_VOUS (Id_Rdv, Dat_Rdv, Heu_Deb, Motif, Statut, #Id_Pat, #Id_Med)\n"
        "CONSULTATION (Id_Cons, Dat_Cons, Diagnostic, #Id_Rdv)\n"
        "LIGNE_ORD (Id_Lig, Dose, Frequence, Duree, #Id_Ord, #Id_Med)\n"
        "ASSURANCE (Id_Assr, Nom_Assr, Type_Contrat, Taux_PEC, #Id_Pat)"
    )
    tf.text = mld_text
    tf.paragraphs[0].font.size = Pt(16)

    prs.save("Presentation_SGRDMS_ACSI_Final.pptx")
    print("Présentation finale générée.")

if __name__ == "__main__":
    create_presentation()
